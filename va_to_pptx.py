# va_to_pptx.py
# Kullanım:
#   python va_to_pptx.py --va va.json --in jama_va.pptx --out outputs/va_filled.pptx

import json, argparse, os, re
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# ---------- helpers ----------
def find_shape_by_name(slide, name):
    if not name: return None
    for shp in slide.shapes:
        if getattr(shp, "name", None) == name:
            return shp
    return None

def set_text(shape, text, font_size=16, bullet=False, shrink_threshold=700):
    if shape is None or not getattr(shape, "has_text_frame", False):
        return
    txt = text or ""
    target_size = Pt(12 if len(txt) > shrink_threshold else font_size)
    tf = shape.text_frame
    tf.clear()
    if not bullet:
        p = tf.paragraphs[0]
        run = p.add_run(); run.text = txt
        run.font.size = target_size; run.font.color.rgb = RGBColor(0,0,0)
    else:
        lines = [ln for ln in (txt.splitlines() or [""]) if ln.strip()!=""] or [""]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
            p.level = 0
            run = p.add_run(); run.text = line
            run.font.size = target_size; run.font.color.rgb = RGBColor(0,0,0)

def ensure_dir(path):
    d = os.path.dirname(path)
    if d: os.makedirs(d, exist_ok=True)

def first_sentence(text: str) -> str:
    t = (text or "").strip()
    if not t: return ""
    return re.split(r"(?<=[.!?])\s+", t, maxsplit=1)[0].strip()

def rest_sentences(text: str) -> str:
    t = (text or "").strip()
    if not t: return ""
    parts = re.split(r"(?<=[.!?])\s+", t, maxsplit=1)
    return parts[1].strip() if len(parts)>1 else ""

def join_nonempty(lines): return "\n".join([s for s in lines if s and s.strip()])

# JSON alanı string/dict olabilir: güvenli çekici
def pick(dct, key, sub=None):
    v = (dct or {}).get(key, "")
    if isinstance(v, dict):
        if sub: return v.get(sub, "")
        return v.get("description","") or v.get("subtitle","")
    return v or ""

# ---------- main ----------
def main(va_path, in_pptx, out_pptx):
    with open(va_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    va = data.get("va", {})
    ts = va.get("the_study", {})
    fd = va.get("findings", {})

    title = data.get("title","")
    url   = data.get("url","")

    # Population
    pop_subtitle = pick(ts, "participants", "subtitle") or first_sentence(pick(ts,"participants"))
    pop_desc     = pick(ts, "participants", "description") or pick(ts,"participants")

    # Intervention
    comparator   = pick(ts, "comparator")
    inter_sub    = pick(ts, "intervention", "subtitle") or (f"Intervention vs {comparator}" if comparator else first_sentence(pick(ts,"intervention")))
    inter_desc   = pick(ts, "intervention", "description") or pick(ts,"intervention")

    # Settings / Locations
    settings_desc = pick(ts, "settings_locations", "description") or pick(ts, "settings_locations")

    # Primary Outcome
    primary_desc  = pick(ts, "primary_outcome", "description") or pick(ts, "primary_outcome")

    # Findings
    # Eğer LLM çıktısı varsa (description_1/2) onu kullan; yoksa summary'den üret.
    f1 = pick(fd, "description_1")
    f2 = pick(fd, "description_2")
    if not (f1 or f2):
        summary = pick(fd, "summary")
        f1 = first_sentence(summary)
        f2 = rest_sentences(summary)

    prs = Presentation(in_pptx)
    slide = prs.slides[0]

    set_text(find_shape_by_name(slide,"title"), title, font_size=22)
    set_text(find_shape_by_name(slide,"footer_citation"), url, font_size=10)

    set_text(find_shape_by_name(slide,"population_subtitle"), pop_subtitle, font_size=16)
    set_text(find_shape_by_name(slide,"population_description"), pop_desc, font_size=14)

    set_text(find_shape_by_name(slide,"intervention_subtitle"), inter_sub, font_size=16)
    set_text(find_shape_by_name(slide,"intervention_description"), inter_desc, font_size=14)

    set_text(find_shape_by_name(slide,"settings_locations_description"), settings_desc, font_size=14)
    set_text(find_shape_by_name(slide,"primary_outcome_description"), primary_desc, font_size=14)

    set_text(find_shape_by_name(slide,"findings_description_1"), f1, font_size=14)
    set_text(find_shape_by_name(slide,"findings_description_2"), f2, font_size=14, bullet=True)

    ensure_dir(out_pptx)
    prs.save(out_pptx)
    print(f"OK -> {out_pptx}")

if __name__ == "__main__":
    ap = argparse.ArgumentParser(description="VA JSON'u JAMA VA şablonuna basar.")
    ap.add_argument("--va", required=True)
    ap.add_argument("--in", dest="in_pptx", required=True)
    ap.add_argument("--out", default="outputs/va_filled.pptx")
    args = ap.parse_args()
    main(args.va, args.in_pptx, args.out)
