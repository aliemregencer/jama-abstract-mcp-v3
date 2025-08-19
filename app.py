import os
import re
import base64
import time
from datetime import datetime
from typing import Optional, Tuple

import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# -------------------- scrape utils --------------------
def _clean(s: Optional[str]) -> str:
    import html, unicodedata
    if not s: return ""
    s = html.unescape(s)
    s = unicodedata.normalize("NFKC", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s.replace("\u2212", "-")

def _norm_heading(h: str) -> Optional[str]:
    h = _clean(h).lower().rstrip(":")
    MAP = {
        "importance":"importance",
        "objective":"objective",
        "design, setting, and participants":"dsp",
        "design, settings, and participants":"dsp",
        "participants":"dsp",
        "intervention":"interventions",
        "interventions":"interventions",
        "main outcomes and measures":"moam",
        "outcomes":"moam",
        "results":"results",
        "conclusions and relevance":"conclusions",
        "conclusions":"conclusions",
        "meaning":"meaning",
        "trial registration":"trial_registration",
    }
    return MAP.get(h)

def _parse_abstract_dom(soup: BeautifulSoup) -> dict:
    out = {}
    abstract = soup.find(id="abstract") or soup.select_one("#abstract")
    if not abstract: return out
    for p in abstract.find_all("p", recursive=True):
        strong = p.find("strong")
        if not strong: continue
        key = _norm_heading(strong.get_text(" ", strip=True))
        if not key: continue
        whole = p.get_text(" ", strip=True)
        content = re.sub(r"^\s*"+re.escape(strong.get_text(strip=True))+r"\s*:?\s*", "", whole, flags=re.I)
        out[key] = _clean(content)
    return out

def _parse_abstract_meta(soup: BeautifulSoup) -> dict:
    out = {}
    meta = soup.find("meta", attrs={"name":"citation_abstract"})
    if not meta or not meta.get("content"): return out
    inner = BeautifulSoup(meta["content"], "lxml")
    for h in inner.find_all(["h3","strong"]):
        key = _norm_heading(h.get_text(" ", strip=True))
        if not key: continue
        content = ""
        sib = h.find_next_sibling()
        if sib and sib.name == "p":
            content = sib.get_text(" ", strip=True)
        else:
            par = h.parent if h.parent and h.parent.name == "p" else None
            if par:
                whole = par.get_text(" ", strip=True)
                content = re.sub(r"^\s*"+re.escape(h.get_text(strip=True))+r"\s*:?\s*", "", whole, flags=re.I)
        if content:
            out[key] = _clean(content)
    return out

def _parse_key_points(soup: BeautifulSoup) -> dict:
    out = {"question":"", "findings":"", "meaning":""}
    hdr = None
    for tag in soup.find_all(["h2","h3","h4"]):
        if _clean(tag.get_text()).lower() == "key points":
            hdr = tag; break
    if not hdr: return out
    container = hdr.parent
    for p in container.find_all("p"):
        t = _clean(p.get_text(" ", strip=True))
        m = re.match(r"(question|findings|meaning)\.?\s*(.+)", t, flags=re.I)
        if m:
            k, v = m.group(1).lower(), _clean(m.group(2))
            out[k] = v
    return out

def _pull_comparator(text: str) -> str:
    t = _clean(text)
    for pat in [r"\bvs\.?\s+([^.;:]+)", r"\bversus\s+([^.;:]+)", r"\bcompared with\s+([^.;:]+)"]:
        m = re.search(pat, t, flags=re.I)
        if m: return _clean(m.group(1))
    return ""

def _pull_settings_locations(sections: dict) -> str:
    dsp = sections.get("dsp","")
    if not dsp: return ""
    loc_rgx = re.compile(
        r"(?:\d+\s+(?:center|centers|site|sites|unit|units|hospital|hospitals)|across the\s+[A-Za-z ,\-]+|in\s+[A-Z][a-zA-Z]+(?:\s+[A-Z][a-zA-Z]+)?|multicenter|single[- ]center)",
        flags=re.I
    )
    sentences = re.split(r"(?<=[.!?])\s+", dsp)
    for s in sentences:
        if loc_rgx.search(s):
            return _clean(s)[:250]
    return ""

def _pull_primary_outcome(moam: str, backup_texts: str = "") -> str:
    for t in [moam, backup_texts]:
        tt = _clean(t)
        m = re.search(r"(primary (?:outcome|endpoint)[^.;:]*[.;:]?)", tt, flags=re.I)
        if m: return _clean(m.group(1))
    return _clean(moam)

def scrape_url(url: str) -> dict:
    r = requests.get(url, headers={"User-Agent":"Mozilla/5.0","Accept-Language":"en-US,en;q=0.9"}, timeout=25)
    r.raise_for_status()
    soup = BeautifulSoup(r.text, "lxml")

    title = ""
    h1 = soup.find("h1")
    if h1: title = _clean(h1.get_text(" ", strip=True))
    if not title:
        og = soup.find("meta", attrs={"property":"og:title"})
        if og and og.get("content"): title = _clean(og["content"])
    if not title:
        ct = soup.find("meta", attrs={"name":"citation_title"})
        if ct and ct.get("content"): title = _clean(ct["content"])

    secs = _parse_abstract_dom(soup) or _parse_abstract_meta(soup)
    kp = _parse_key_points(soup)

    participants = secs.get("dsp","")
    intervention = secs.get("interventions","")
    moam_text    = secs.get("moam","")
    results      = secs.get("results","")
    importance   = secs.get("importance","")
    objective    = secs.get("objective","")
    conclusions  = secs.get("conclusions","") or secs.get("meaning","")

    before       = kp["question"]  or importance
    findings_sum = kp["findings"]  or results
    implications = kp["meaning"]   or conclusions

    comparator   = _pull_comparator(intervention) or _pull_comparator(participants)
    settings_locs= _pull_settings_locations(secs)
    primary_out  = _pull_primary_outcome(moam_text, participants + " " + intervention)

    return {
        "url": url,
        "title": title,
        "va": {
            "the_study": {
                "participants": participants,
                "intervention": intervention,
                "comparator": comparator,
                "primary_outcome": primary_out,
                "settings_locations": settings_locs
            },
            "findings": {
                "summary": findings_sum,
                "key_numbers": []
            }
        }
    }

# -------------------- PPTX render --------------------
def _find_shape(slide, name):
    for shp in slide.shapes:
        if getattr(shp, "name", None) == name:
            return shp
    return None

def _set_text(shape, txt, size=16):
    if shape is None or not getattr(shape, "has_text_frame", False): return
    tf = shape.text_frame; tf.clear()
    run = tf.paragraphs[0].add_run()
    run.text = txt or ""
    run.font.size = Pt(size); run.font.color.rgb = RGBColor(0,0,0)

def _first_sentence(t: str) -> str:
    t = (t or "").strip()
    if not t: return ""
    return re.split(r"(?<=[.!?])\s+", t, maxsplit=1)[0].strip()

def render_to_pptx(data: dict, template_path: str, output_path: str) -> str:
    va = data.get("va", {})
    ts = va.get("the_study", {})
    fd = va.get("findings", {})

    prs = Presentation(template_path)
    slide = prs.slides[0]

    _set_text(_find_shape(slide,"title"), data.get("title",""), 22)
    _set_text(_find_shape(slide,"footer_citation"), data.get("url",""), 10)

    pop = ts.get("participants","")
    _set_text(_find_shape(slide,"population_subtitle"), _first_sentence(pop), 16)
    _set_text(_find_shape(slide,"population_description"), pop, 14)

    inter = ts.get("intervention","")
    comp  = ts.get("comparator","")
    inter_sub = f"Intervention vs {comp}" if comp else _first_sentence(inter)
    _set_text(_find_shape(slide,"intervention_subtitle"), inter_sub, 16)
    _set_text(_find_shape(slide,"intervention_description"), inter, 14)

    _set_text(_find_shape(slide,"settings_locations_description"), ts.get("settings_locations",""), 14)
    _set_text(_find_shape(slide,"primary_outcome_description"), ts.get("primary_outcome",""), 14)

    summary = fd.get("summary","")
    _set_text(_find_shape(slide,"findings_description_1"), _first_sentence(summary), 14)
    rest = re.split(r"(?<=[.!?])\s+", (summary or "").strip(), maxsplit=1)
    _set_text(_find_shape(slide,"findings_description_2"), rest[1].strip() if len(rest)>1 else "", 14)

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    return output_path

# -------------------- GitHub upload --------------------
def upload_to_github_release(
    filename: str,
    title: str,
    repo_full_name: str,
    github_token: str,
) -> Tuple[Optional[str], Optional[str]]:
    """
    'latest-abstract' tag'li release yaratır (varsa silip baştan), PPTX'i asset olarak yükler,
    herkese açık browser_download_url döndürür.
    """
    try:
        if not repo_full_name or "/" not in repo_full_name:
            return None, "Geçersiz repo formatı. 'kullanici/repoadi' olmalı."
        if not github_token:
            return None, "GitHub token gerekli."

        owner, repo = repo_full_name.split("/", 1)
        api_base = f"https://api.github.com/repos/{owner}/{repo}"
        headers_json = {
            "Authorization": f"Bearer {github_token}",
            "Accept": "application/vnd.github+json",
        }
        tag = "latest-abstract"
        safe_title = (title or "JAMA Abstract")[:70]

        # repo erişimi
        repo_check = requests.get(api_base, headers=headers_json)
        if repo_check.status_code != 200:
            return None, f"Repo erişimi başarısız: {repo_check.status_code} {repo_check.text}"

        # eski release + tag sil
        r = requests.get(f"{api_base}/releases/tags/{tag}", headers=headers_json)
        if r.status_code == 200:
            rel = r.json()
            rid = rel["id"]
            assets = requests.get(f"{api_base}/releases/{rid}/assets", headers=headers_json).json()
            for a in assets:
                requests.delete(f"{api_base}/releases/assets/{a['id']}", headers=headers_json)
            requests.delete(f"{api_base}/releases/{rid}", headers=headers_json)
            requests.delete(f"{api_base}/git/refs/tags/{tag}", headers=headers_json)

        # yeni release
        rel_body = {
            "tag_name": tag,
            "name": f"JAMA Abstract - {safe_title}",
            "body": f"Otomatik üretilmiş görsel özet\nMakale: {safe_title}\nTarih: {datetime.now():%Y-%m-%d %H:%M:%S}",
            "draft": False,
            "prerelease": False
        }
        cr = requests.post(f"{api_base}/releases", json=rel_body, headers=headers_json)
        # özel durum: boş repo
        if cr.status_code == 422 and "Repository is empty" in cr.text:
            # boş repoyu README ile başlat
            default_branch = (requests.get(api_base, headers=headers_json).json().get("default_branch")) or "main"
            readme = "# Auto Init\n\nPPTX assets for visual abstracts."
            init = requests.put(
                f"{api_base}/contents/README.md",
                headers=headers_json,
                json={"message":"init", "content": base64.b64encode(readme.encode()).decode(), "branch": default_branch}
            )
            if init.status_code not in (200,201):
                return None, f"README oluşturulamadı: {init.status_code} {init.text}"
            time.sleep(1)
            cr = requests.post(f"{api_base}/releases", json=rel_body, headers=headers_json)

        if cr.status_code != 201:
            return None, f"Release oluşturma hatası: {cr.status_code} {cr.text}"

        rel = cr.json()
        upload_url = rel["upload_url"].split("{")[0]

        with open(filename, "rb") as f:
            binary = f.read()
        headers_upload = {
            "Authorization": f"Bearer {github_token}",
            "Accept": "application/vnd.github+json",
            "Content-Type": "application/octet-stream",
        }
        ur = requests.post(f"{upload_url}?name={os.path.basename(filename)}", data=binary, headers=headers_upload)
        if ur.status_code != 201:
            return None, f"Dosya yükleme hatası: {ur.status_code} {ur.text}"

        asset = ur.json()
        return asset.get("browser_download_url"), None
    except Exception as e:
        return None, f"GitHub yükleme hatası: {e}"

# -------------------- Public API --------------------
def create_graphical_abstract_from_url(url: str) -> Tuple[str, str]:
    """
    Yalnızca yerelde PPTX üretir. MCP `server.py` bunu döndürür.
    """
    data = scrape_url(url)

    # (İsteğe bağlı) LLM kısaltma adımı: burada data['va'] üzerinde uygulanabilir.
    # if os.getenv("USE_LLM_SHORTEN") == "1": data = run_llm_shorten(data)

    template = os.environ.get("JAMA_TEMPLATE", "templates/abstract.pptx")
    out_dir  = os.environ.get("OUTPUT_DIR", "outputs")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "visual_abstract.pptx")
    render_to_pptx(data, template, out_path)
    return out_path, "PPTX oluşturuldu."

def create_graphical_abstract(url: str, github_repo: str, github_token: str) -> Tuple[str, Optional[str], str]:
    """
    PPTX üret + GitHub release'e yükle. (output_path, download_url, mesaj)
    """
    out_path, _ = create_graphical_abstract_from_url(url)
    # data.title'ı yeniden almak için küçük bir istek (alternatif: render aşamasında başlığı da dön)
    title = os.path.basename(out_path)
    # daha iyi başlık için sayfadan tekrar oku:
    try:
        title = scrape_url(url).get("title") or title
    except Exception:
        pass

    download_url, err = upload_to_github_release(out_path, title, github_repo, github_token)
    if download_url:
        return out_path, download_url, "PPTX oluşturuldu ve release'e yüklendi."
    return out_path, None, f"PPTX oluşturuldu, fakat yükleme başarısız: {err}"
